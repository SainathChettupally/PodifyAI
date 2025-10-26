from __future__ import annotations
from pathlib import Path
import re

def extract_text_from_pdf(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    pages = []
    for p in doc:
        pages.append(p.get_text("text"))
    return "\n".join(pages)

def extract_text_from_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    return "\n".join([p.text for p in d.paragraphs])

def extract_text_from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(" ".join(run.text for run in para.runs))
    return "\n".join(t for t in texts if t and t.strip())

def extract_text_from_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def extract_text_from_html(path: str) -> str:
    from bs4 import BeautifulSoup
    html = Path(path).read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(" ")

def extract_text_from_csv(path: str, max_rows: int = 5000) -> str:
    import csv
    rows = []
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i > max_rows:
                break
            rows.append(" ".join(row))
    return "\n".join(rows)

def clean_text(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()

def truncate_for_demo(s: str, max_chars: int = 5000) -> str:
    return s[:max_chars]

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".html", ".htm", ".csv"}

def detect_and_extract_text(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext == ".docx":
        return extract_text_from_docx(path)
    if ext == ".pptx":
        return extract_text_from_pptx(path)
    if ext in {".txt", ".md"}:
        return extract_text_from_txt(path)
    if ext in {".html", ".htm"}:
        return extract_text_from_html(path)
    if ext == ".csv":
        return extract_text_from_csv(path)
    raise ValueError(f"Unsupported file type: {ext}. Supported: {sorted(SUPPORTED_EXTS)}")
